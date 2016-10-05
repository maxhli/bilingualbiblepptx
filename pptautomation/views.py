from django.shortcuts import render

from django.core.files import File
from io import BytesIO
from os.path import (
    sep,
    getsize)

from os.path import dirname
from django.conf import settings
from wsgiref.util import FileWrapper
# Create your views here.
from pptx import Presentation
from pptx.util import Inches, Pt
# from pptx.enum.shapes import MSO_SHAPE
# from pptx.enum.dml import MSO_THEME_COLOR
# from pptx.enum.text import PP_ALIGN
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

# from pptx.dml.color import RGBColor
from django.http import (
    HttpResponse,
    FileResponse)


def main(request):
    return HttpResponse("I am good.")


def some_view(request):

    # with open('test.txt') as f:
    #    paragraph_strs = f.readlines()
    # paragraph_strs = [line.rstrip('\n') for line in open('test.txt')]

    paragraph_strs = [
        '1 〔大衛的詩、交與伶長。〕諸天述說　神的榮耀．穹蒼傳揚他的手段。',
        '2 這日到那日發出言語．這夜到那夜傳出知識。',
        '1 For the director of music. A psalm of David. The heavens declare the glory of God; the skies proclaim the work of his hands.',
        '2 Day after day they pour forth speech; night after night they display knowledge.'
    ]


    background = settings.BASE_DIR + sep + 'starting_background.pptx'
    print(background)
    prs = Presentation(background)
    # prs = Presentation()
    # # prs = Presentation()
    #
    # # default slide width
    # # prs.slide_width = 9144000
    # # slide height @ 4:3
    # # prs.slide_height = 6858000
    # # slide height @ 16:9
    # prs.slide_height = 5143500
    #
    # slide_master = prs.slide_masters[0]
    # slide_layout = prs.slide_layouts[1]
    #
    # #slide = prs.slides[0]
    #
    # SLD_LAYOUT_TITLE_AND_CONTENT = 1
    #
    # slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
    # slide = prs.slides.add_slide(slide_layout)
    #
    # for shape in slide.shapes:
    #     if not shape.has_text_frame:
    #         continue
    #     text_frame = shape.text_frame
    #
    # text_frame.clear()  # remove any existing paragraphs, leaving one empty one
    #
    # # p = text_frame.paragraphs[0]
    # # p.text = paragraph_strs[0]
    #
    # for para_str in paragraph_strs:
    #     print(len(para_str), " ", para_str)
    #     p = text_frame.add_paragraph()
    #     p.text = para_str
    #     p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    #     p.font.size = Pt(23)

    # prs = Presentation(buffer)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = Inches(0.5)

    width = Inches(9)
    height = Inches(4)

    tx_box = slide.shapes.add_textbox(left, top, width, height)

    tf = tx_box.text_frame

    for para_str in paragraph_strs:
        print(len(para_str), " ", para_str)
        p = tf.add_paragraph()

        p.text = para_str
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.font.size = Pt(24)
        p.level = 0

    prs.save('oldgen.pptx')

    filename = settings.BASE_DIR + sep + 'oldgen.pptx'
    file_content = open(filename, "rb")
    response = HttpResponse(
        file_content,
        content_type=
        'application/vnd.openxmlformats-officedocument.presentationml.presentation')

    response['Content-Disposition'] = 'attachment; filename="oldgen.pptx"'
    response['Content-Length'] = getsize(filename)

    print("File size is: " + str(getsize(filename)))
    return response
